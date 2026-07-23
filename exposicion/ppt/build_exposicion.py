#!/usr/bin/env python3
"""Genera Exposicion_Plan_Tesis_Grupo9.2.pptx según ppt-rules.pptx (Grupo 9.2)."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent
ASSETS = ROOT / "assets"
TEMPLATE = ROOT / "ppt-rules.pptx"
OUTPUT = ROOT / "Exposicion_Plan_Tesis_Grupo9.2.pptx"

FONT = "Calibri"
BODY = Pt(22)
BODY_SM = Pt(20)
TABLE = Pt(14)
CAPTION = Pt(18)

MARGIN_X = Inches(0.55)
# Imágenes: ancho fijo, sin forzar alto (evita recorte/distorsión). Ajuste fino manual en PowerPoint.
IMG_LEFT = Inches(0.6)
IMG_TOP = Inches(1.35)
IMG_WIDTH = Inches(12.1)
CAPTION_TOP = Inches(6.85)
FOOTER = Inches(0.28)


def delete_slide(prs: Presentation, index: int) -> None:
    slide_id = prs.slides._sldIdLst[index]
    prs.part.drop_rel(slide_id.rId)
    del prs.slides._sldIdLst[index]


def clear_slides(prs: Presentation) -> None:
    for i in range(len(prs.slides) - 1, -1, -1):
        delete_slide(prs, i)


def set_notes(slide, text: str) -> None:
    slide.notes_slide.notes_text_frame.text = text


def add_title_image(
    prs,
    title: str,
    image_name: str,
    caption: str = "",
    notes: str = "",
):
    """Inserta figura con ancho fijo; PowerPoint conserva proporción (sin recortar)."""
    slide = prs.slides.add_slide(prs.slide_layouts[4])
    slide.shapes.title.text = title
    path = ASSETS / image_name
    if path.exists():
        slide.shapes.add_picture(str(path), IMG_LEFT, IMG_TOP, width=IMG_WIDTH)
    if caption:
        box = slide.shapes.add_textbox(MARGIN_X, CAPTION_TOP, prs.slide_width - 2 * MARGIN_X, Inches(0.55))
        p = box.text_frame.paragraphs[0]
        p.text = caption
        p.font.name = FONT
        p.font.size = CAPTION
        p.alignment = PP_ALIGN.CENTER
    if notes:
        set_notes(slide, notes)
    return slide


def add_bullets(prs, title: str, bullets: list[str], notes: str = "", size=BODY):
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.name = FONT
        p.font.size = size
    if notes:
        set_notes(slide, notes)
    return slide


def add_two_cols(prs, title: str, left: list[str], right: list[str], notes: str = ""):
    slide = prs.slides.add_slide(prs.slide_layouts[3])
    slide.shapes.title.text = title
    for idx, lines in enumerate((left, right)):
        tf = slide.placeholders[1 + idx].text_frame
        tf.clear()
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line
            p.font.name = FONT
            p.font.size = BODY_SM
    if notes:
        set_notes(slide, notes)
    return slide


def add_table_slide(
    prs,
    title: str,
    headers: list[str],
    rows: list[list[str]],
    notes: str = "",
    col_widths: list[float] | None = None,
):
    slide = prs.slides.add_slide(prs.slide_layouts[4])
    slide.shapes.title.text = title
    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_top = Inches(1.15)
    table_left = MARGIN_X
    table_width = prs.slide_width - 2 * MARGIN_X
    table_height = prs.slide_height - table_top - FOOTER - Inches(0.15)
    shape = slide.shapes.add_table(n_rows, n_cols, table_left, table_top, table_width, table_height)
    table = shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = int(table_width * w)
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        for p in cell.text_frame.paragraphs:
            p.font.name = FONT
            p.font.size = TABLE
            p.font.bold = True
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = val
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for p in cell.text_frame.paragraphs:
                p.font.name = FONT
                p.font.size = TABLE
    if notes:
        set_notes(slide, notes)
    return slide


def build():
    prs = Presentation(str(TEMPLATE))
    clear_slides(prs)

    # --- ppt-rules slide 1: Título ---
    s1 = prs.slides.add_slide(prs.slide_layouts[0])
    s1.shapes.title.text = (
        "Alineación contratos OpenAPI bancarios\n↔ Service Domains BIAN"
    )
    sub = s1.placeholders[1].text_frame
    sub.clear()
    for i, line in enumerate(
        [
            "Investigación tecnológica · Design Science",
            "IEEE: Knowledge Representation · ACM: Semantic networks",
            "Grupo 9.2 — Alex Mancilla · Jack Paitan",
            "Universidad Nacional de Ingeniería · Julio 2026",
        ]
    ):
        p = sub.paragraphs[0] if i == 0 else sub.add_paragraph()
        p.text = line
        p.font.name = FONT
        p.font.size = BODY
    set_notes(s1, "[0:00–0:40] Jack. Título + tipo investigación + clasificación ACM/IEEE.")

    # --- ppt-rules slide 4: Introducción (1 transparencia) ---
    add_bullets(
        prs,
        "Introducción al tema (§1.1)",
        [
            "OpenAPI: contrato YAML/JSON machine-readable (Casas et al., 2021)",
            "BIAN: Service Domains como referencia de negocio (Farzi, 2021)",
            "Alineación = correspondencia estructural + semántica contrato ↔ SD",
            "Producto del plan: S, E, C → AlignmentScore (procedimiento reproducible)",
            "Guía 1: diseño documentado; validación empírica → Tesis 2 (§1.6)",
        ],
        notes="[0:40–1:30] Alex. Conceptos antes del problema (PI-12).",
    )

    # --- ppt-rules slide 5: Ontología (1 transparencia) ---
    add_title_image(
        prs,
        "Ontología del plan — conceptos relacionados",
        "Figura_02_Objeto_Unidad_Analisis.png",
        "Contrato OpenAPI · UA instancia · Producto S/E/C · Referencia BIAN",
        notes="[1:30–2:15] Jack. Señalar cajas; ID {fuente}_{SD}_{versión}.",
    )

    # --- ppt-rules slides 6–7: Problemática real ---
    add_bullets(
        prs,
        "Problemática — situación real: dónde y actores (§1.2)",
        [
            "Dónde: ecosistema bancario digital — APIs propias, fintechs, open banking (CMA, 2016)",
            "Ambiente virtual: contratos OpenAPI publicados antes del despliegue en canales",
            "Actores: arquitectura, integración, gobernanza de APIs, socios externos",
            "Proceso administrativo: revisión y certificación de interfaces REST",
            "La problemática no se resuelve directamente con el artefacto del plan",
        ],
        notes="[2:15–3:00] Jack. Situación real ≠ objeto ≠ problema tecnológico.",
    )

    add_bullets(
        prs,
        "Problemática — síntomas, dolor y estadísticas (§1.2)",
        [
            "Síntomas: fragmentación semántica, costos de integración, deriva negocio–REST",
            "Dolor: retrabajo en adaptadores, ensayo–error, riesgo en releases (TSB, 2018)",
            "Casas et al. (2021): 68 % papers proponen herramientas; 43 % documentación",
            "Actores más afectados: equipos de integración y gobernanza de APIs",
            "Motiva medir coherencia documental antes de exponer servicios",
        ],
        notes="[3:00–3:45] Jack. Tener PDF Casas listo (PI-13).",
    )

    # --- ppt-rules slide 8: Objeto de estudio (subyacente) ---
    add_bullets(
        prs,
        "Objeto de estudio — situación subyacente (§1.4)",
        [
            "Causa de síntomas: desalineación entre capa REST y modelo de negocio BIAN",
            "Objeto: contrato OpenAPI bancario completo (archivo YAML/JSON)",
            "Características: paths, operations, schemas, metadatos machine-readable",
            "Dos jerarquías comparables: OpenAPI ↔ Service Domain BIAN",
            "Usar ontología (slide anterior) para exponer el problema subyacente",
        ],
        notes="[3:45–4:30] Jack. BIAN = referencia, no objeto.",
    )

    # --- ppt-rules slide 9: Objeto real ---
    add_bullets(
        prs,
        "Objeto de estudio real — taxonomía y ciclo de vida",
        [
            "Entidad real: contrato OpenAPI publicado por un banco (independiente del observador)",
            "Taxonomía: OpenAPI 3.x · dominios pagos / cuentas / préstamos · versiones",
            "Ciclo de vida: diseño → revisión → publicación → consumo (no tráfico runtime)",
            "Registro: archivo versionado en repositorio de APIs / portal de desarrolladores",
            "Ejemplos Anexo F: pares contrato ↔ Payment Initiation SD",
        ],
        notes="[4:30–5:00] Jack. Ciclo de vida del contrato, no del proceso bancario completo.",
    )

    # --- ppt-rules slides 10–11: Instrumento + objeto modelado ---
    add_two_cols(
        prs,
        "Instrumento y objeto modelado",
        [
            "Instrumento de captura:",
            "• Parser OpenAPI 3.x + extracto BIAN",
            "• Procedimiento de normalización (OE1)",
            "• Schema matching + similitud semántica",
            "Objeto modelado:",
            "Representaciones normalizadas + matrices + scores S/E/C",
        ],
        [
            "Ejemplos de instancias (Anexo F):",
            "Caso A: POST /payments vs InitiatePayment → afecta S",
            "Caso B: GET /payments/{id} vs 3 behaviors → afecta C",
            "Estados de clasificación del artefacto:",
            "Alta ≥0,80 · Media 0,60–0,79",
            "Baja 0,40–0,59 · Nula <0,40",
        ],
        notes="[5:00–5:45] Alex. Instrumento = procedimiento, no encuesta ni IoT.",
    )

    # --- ppt-rules slides 12–14: Problema tecnológico + PE ---
    add_bullets(
        prs,
        "Problema tecnológico y preguntas P1–P4 (§1.3)",
        [
            "Artefacto: procedimiento + modelo S/E/C → AlignmentScore (software documentado)",
            "Tipo: clasificación + regresión acotada [0,1] por score",
            "Problema: métodos insuficientes / no adaptados a OpenAPI↔BIAN (Shvaiko, 2005; Casas, 2021)",
            "P1 matching estructural · P2 similitud semántica · P3 cobertura · P4 score integrado",
            "Solo relacionado al artefacto a diseñar — no a la problemática del sector",
        ],
        notes="[5:45–6:30] Alex. «Insuficiente», no «no existe» (PI-12).",
    )

    # --- ppt-rules slide 15: Objetivo general ---
    add_bullets(
        prs,
        "Objetivo general (técnico) — insumo y variable dependiente",
        [
            "Insumo inicial: contrato OpenAPI bancario + Service Domain BIAN emparejado",
            "OG: diseñar modelo S/E/C y procedimiento reproducible de alineación",
            "Variable dependiente del artefacto: AlignmentScore ∈ [0,1] + clasificación",
            "Objetivo superior (consecuencia): cuantificar coherencia antes del despliegue",
            "Medida de desempeño: S, E, C y score integrado por contrato × SD",
        ],
        notes="[6:30–7:00] Alex. Separar OG técnico vs objetivo superior.",
    )

    # --- ppt-rules slide 16: OE específicos ---
    add_title_image(
        prs,
        "Objetivos específicos OE1–OE4 (entrada → producto)",
        "Figura_04_Cadena_Objetivos.png",
        "Cada OE es hito con producto medible — no «programar» ni «experimentar»",
        notes="[7:00–7:45] Alex. OE1 normaliza · OE2→E · OE3→S · OE4→C + AlignmentScore.",
    )

    # --- ppt-rules slides 17–18: Variables independientes ---
    add_table_slide(
        prs,
        "Variables independientes — factores del investigador (Tabla 5)",
        ["OE", "Factor (VI)", "Estados / alternativas"],
        [
            ["OE1", "Esquema de representación intermedia", "Paths/schemas vs behaviors/BO normalizados"],
            ["OE2", "Estrategia de schema matching", "Correspondencia estructural + umbral fijo"],
            ["OE3", "Técnica de similitud semántica", "Embeddings / ontología / híbrido"],
            ["OE4", "Pesos α, β, γ de agregación", "α+β+γ=1; taxonomía de tipologías"],
        ],
        notes="[7:45–8:15] Alex. VI = factores de diseño, no hiperparámetros de entrenamiento.",
        col_widths=[0.12, 0.38, 0.50],
    )

    # --- ppt-rules slide 19: Variables dependientes ---
    add_table_slide(
        prs,
        "Variables dependientes — ratios por OE (Tabla 5 · Anexo D)",
        ["OE", "Ratio (VD)", "Rango / fórmula"],
        [
            ["OE1", "Completitud de normalización", "Elementos extraídos / esperados"],
            ["OE2", "StructuralScore (E)", "E ∈ [0, 1] — schema matching"],
            ["OE3", "SemanticScore (S)", "S ∈ [0, 1] — similitud semántica"],
            ["OE4", "CoverageScore (C); AlignmentScore", "C ∈ [0,1]; Score = αS+βE+γC"],
        ],
        notes="[8:15–8:45] Alex. VD del artefacto = AlignmentScore + clasificación.",
        col_widths=[0.10, 0.42, 0.48],
    )

    # --- ppt-rules slide 20: Metodología ---
    add_bullets(
        prs,
        "Metodología — hitos / fases con verbos (§4.2 · §4.2.1)",
        [
            "Fase 1 · Diseñar entradas, scores y umbrales → modelo S/E/C documentado",
            "Fase 2 · Normalizar contrato + SD → representaciones comparables (OE1)",
            "Fase 3 · Emparejar estructuras → StructuralScore E (OE2)",
            "Fase 4 · Calcular similitud → SemanticScore S (OE3)",
            "Fase 5 · Integrar S,E,C → CoverageScore + AlignmentScore + tipologías (OE4)",
        ],
        notes="[8:45–9:15] Alex. Partir por metodología (PI-14 cangrejo).",
    )

    # --- Artefacto: modelo S/E/C (ppt-rules slide 12 arquitectura) ---
    add_title_image(
        prs,
        "Artefacto — modelo de alineación S / E / C",
        "Figura_03_AlignmentScore.png",
        "AlignmentScore = α·S + β·E + γ·C  →  Alta / Media / Baja / Nula",
        notes="[9:15–9:45] Alex. Shvaiko (2005) sustenta E; scores por UA, agregado por contrato×SD.",
    )

    # --- ppt-rules slide 21: Matriz de consistencia (I) ---
    add_table_slide(
        prs,
        "Matriz de consistencia (I) — Anexo C",
        ["Fase", "Problema tecnológico", "Objetivo", "Producto verificable"],
        [
            [
                "Fase 1",
                "Método insuficiente para cuantificar alineación OpenAPI↔BIAN",
                "Diseñar modelo S/E/C y procedimiento",
                "Modelo metodológico documentado (Anexo C)",
            ],
            [
                "Fase 2",
                "Representaciones no comparables entre contrato y SD",
                "OE1 · Normalizar artefactos",
                "Representaciones normalizadas versionadas",
            ],
            [
                "Fase 3",
                "Sin StructuralScore sistemático",
                "OE2 · Emparejar estructuras",
                "Matriz correspondencias + E ∈ [0,1]",
            ],
        ],
        notes="[9:45–10:15] Alex. Matriz ≠ brecha Cap. 3. Problema tecnológico ≠ problemática §1.2.",
        col_widths=[0.10, 0.30, 0.28, 0.32],
    )

    # --- Matriz de consistencia (II) ---
    add_table_slide(
        prs,
        "Matriz de consistencia (II) — VI · VD · Hipótesis",
        ["Fase", "Variable independiente", "Variable dependiente", "Hipótesis (síntesis)"],
        [
            ["Fase 1", "Esquema S/E/C; pesos α,β,γ", "Definición scores", "H0: modelo S/E/C reproducible"],
            ["Fase 2", "Esquema intermedio", "Completitud normalización", "H1 (OE1)"],
            ["Fase 3", "Schema matching", "StructuralScore E", "H2 (OE2)"],
            ["Fase 4", "Similitud semántica", "SemanticScore S", "H3 (OE3)"],
            ["Fase 5", "Pesos α,β,γ", "C; AlignmentScore; clasif.", "H4 (OE4)"],
        ],
        notes="[10:15–10:45] Alex. Señalar fila Fase 1–OE4; matriz completa en PDF Anexo C.",
        col_widths=[0.10, 0.24, 0.28, 0.38],
    )

    # --- Brecha (complemento plan; distinguir de matriz) ---
    add_title_image(
        prs,
        "Brecha de investigación V1–V5 (§3.6) — no confundir con matriz",
        "Figura_06_Brecha.png",
        "Brecha = estado del arte · Matriz = problema tecnológico del diseño",
        notes="[10:45–11:05] Jack: V1–V2 · Alex: V3–V5.",
    )

    # --- ppt-rules slide 22: Conclusiones ---
    add_bullets(
        prs,
        "Conclusiones — relaciones entre conceptos",
        [
            "OpenAPI maduro técnicamente; falta procedimiento OpenAPI↔BIAN reproducible",
            "Objeto (contrato) + referencia (BIAN) + UA (instancia) → scores S/E/C integrados",
            "OE1–OE4 encadenan transformaciones con productos verificables",
            "Matriz Anexo C alinea problema tecnológico, objetivos, VI/VD e hipótesis",
            "Guía 1 entrega diseño; Tesis 2: artefacto software y validación empírica",
        ],
        notes="[11:05–11:25] Jack. Concluir con relaciones, no opinión (ppt-rules).",
    )

    # --- ppt-rules slide 23: Recomendaciones ---
    add_bullets(
        prs,
        "Recomendaciones",
        [
            "Implementar artefacto software según procedimiento §4.2 (Tesis 2)",
            "Calibrar umbrales 0,80/0,60/0,40 con muestra bancaria representativa",
            "Validar hipótesis H1–H4 con contratos OpenAPI 3.x públicos (Anexo F)",
            "Integrar score en pipeline de gobernanza de APIs antes del despliegue",
        ],
        notes="[11:25–11:40] Alex. Breve; foco en continuidad post-Guía 1.",
        size=BODY_SM,
    )

    # --- ppt-rules slide 24: Referencias (mencionadas en exposición) ---
    add_bullets(
        prs,
        "Referencias citadas en esta exposición",
        [
            "Casas, L., Pinto, S. & Silva, J. (2021). OpenAPI specification survey. J. Syst. Softw.",
            "Shvaiko, P. & Euzenat, J. (2005). A survey of schema-based matching approaches. J. Data Semant.",
            "Farzi, H. et al. (2021). BIAN service landscape. BIAN documentation.",
            "CMA (2016). Retail banking market investigation. UK Competition and Markets Authority.",
            "UK Treasury / TSB (2018). Independent review of TSB IT failure.",
        ],
        notes="[11:40–12:00] Jack. PDFs Casas y Shvaiko abiertos para preguntas. ¿Preguntas?",
        size=BODY_SM,
    )

    prs.save(str(OUTPUT))
    print(f"Generado: {OUTPUT} ({len(prs.slides)} diapositivas)")


if __name__ == "__main__":
    build()
