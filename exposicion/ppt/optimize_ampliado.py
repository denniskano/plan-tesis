#!/usr/bin/env python3
"""
LEGACY — Obsoleto. Presentación oral oficial: ../Exposicion_Plan_Tesis_Grupo9.2-beamer.pdf
Guion: ../LIBRETO_EXPOSICION_COMPLETO.md · GUION_15MIN_AMPLIADO.md

Optimiza Exposicion_Plan_Tesis_Grupo9.2-ampliado.pptx (versión anterior, no usar en clase):
- 15 min de exposición (notas orales por slide)
- ppt-rules: Fuentes NO se expone; resto según guion
- Elimina duplicados de imágenes Alex (misma sección)
- Corrige título portada
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent
INPUT = ROOT / "Exposicion_Plan_Tesis_Grupo9.2-ampliado.pptx"
OUTPUT = ROOT / "Exposicion_Plan_Tesis_Grupo9.2-ampliado.pptx"
BACKUP = ROOT / "Exposicion_Plan_Tesis_Grupo9.2-ampliado.backup.pptx"

FONT = "Calibri"

# Diapositivas a eliminar (1-based): duplicados de problemática/objeto/instrumento
DELETE_SLIDES_1BASED = [6, 9, 10, 14, 15]

# Guion 15 min — (minuto inicio, minuto fin, quien, nota oral)
# Tras borrar 5 slides, el orden expuesto queda en 23 slides (slide 2 Fuentes = NO EXPONER)
TIMING: list[tuple[str, str, str, str]] = [
    ("0:00", "0:45", "Jack", "Título + investigación tecnológica + IEEE/ACM. Señalar slide."),
    ("—", "—", "—", "SLIDE 2 FUENTES: NO EXPONER (ppt-rules §3). Solo entrega Drive."),
    ("0:45", "1:35", "Alex", "Intro: OpenAPI, BIAN, alineación, producto S/E/C. Bullets, no leer párrafo."),
    ("1:35", "2:20", "Jack", "Ontología: señalar objeto, UA, producto, referencia BIAN."),
    ("2:20", "3:10", "Jack", "Problemática situación real: figura + dónde/actores. ≠ problema técnico."),
    ("3:10", "4:05", "Jack", "Síntomas TSB/open banking + Casas 68%. Señalar ambas figuras."),
    ("4:05", "4:50", "Jack", "Objeto subyacente: contrato OpenAPI vs SD BIAN. Figura §1.4."),
    ("4:50", "5:30", "Jack", "Objeto real: taxonomía OpenAPI 3.x y ciclo de vida del contrato."),
    ("5:30", "6:15", "Alex", "Instrumento: parser + normalización. Figura procedimiento."),
    ("6:15", "7:00", "Alex", "Casos Anexo F / estados Alta–Nula. Caso A (S) y B (C)."),
    ("7:00", "7:50", "Alex", "Problema tecnológico + P1–P4. Decir «insuficiente», no «no existe»."),
    ("7:50", "8:25", "Alex", "OG: insumo contrato+SD; VD = AlignmentScore."),
    ("8:25", "9:15", "Alex", "OE1–OE4 cadena. Producto medible por objetivo."),
    ("9:15", "9:55", "Alex", "Factores VI (Tabla 5). No confundir con hiperparámetros ML."),
    ("9:55", "10:35", "Alex", "Ratios VD: S, E, C, AlignmentScore. Fórmula αS+βE+γC."),
    ("10:35", "11:15", "Alex", "Metodología: 5 fases con verbos (§4.2.1). Partir aquí si preguntan."),
    ("11:15", "12:00", "Alex", "Artefacto S/E/C. Shvaiko sustenta E. Score por UA."),
    ("12:00", "12:40", "Alex", "Matriz I: Problema→Objetivo→Producto (Fases 1–3). ≠ brecha."),
    ("12:40", "13:20", "Alex", "Matriz II: VI, VD, H1–H4. Anexo C completo en PDF."),
    ("13:20", "14:00", "Ambos", "Brecha V1–V5. Jack V1–V2; Alex V3–V5. ≠ matriz."),
    ("14:00", "14:30", "Jack", "Conclusiones: relaciones concepto→artefacto. Sin opinión."),
    ("14:30", "14:45", "Alex", "Recomendaciones Tesis 2 (breve)."),
    ("14:45", "15:00", "Jack", "Referencias Casas + Shvaiko. ¿Preguntas? PDFs abiertos."),
]


def delete_slide(prs: Presentation, index: int) -> None:
    slide_id = prs.slides._sldIdLst[index]
    prs.part.drop_rel(slide_id.rId)
    del prs.slides._sldIdLst[index]


def set_notes(slide, header: str, body: str = "") -> None:
    text = header if not body else f"{header}\n\n{body}"
    slide.notes_slide.notes_text_frame.text = text


def fix_title_slide(slide) -> None:
    if not slide.shapes.title:
        return
    slide.shapes.title.text = (
        "Alineación entre contratos OpenAPI bancarios\n"
        "↔ Service Domains BIAN · modelo S/E/C"
    )
    for shape in slide.shapes:
        if shape.has_text_frame and shape != slide.shapes.title:
            tf = shape.text_frame
            tf.clear()
            lines = [
                "Investigación tecnológica · Design Science",
                "IEEE: Knowledge Representation · ACM: Semantic networks",
                "Grupo 9.2 — Alex Mancilla · Jack Paitan · UNI · Julio 2026",
            ]
            for i, line in enumerate(lines):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = line
                p.font.name = FONT
                p.font.size = Pt(22)


def fix_fuentes_slide(slide) -> None:
    if slide.shapes.title:
        slide.shapes.title.text = "Fuentes — NO EXPONER (solo entrega Drive)"
    for shape in slide.shapes:
        if shape.has_text_frame and shape != slide.shapes.title:
            tf = shape.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = (
                "ppt-rules §3: esta diapositiva NO se presenta en clase.\n"
                "Entregar link Drive + PDF informe + Zotero al docente."
            )
            p.font.name = FONT
            p.font.size = Pt(20)
    set_notes(
        slide,
        "⛔ NO EXPONER — omitir en los 15 min.",
        "Avanzar directo de slide 1 → slide 3 (Introducción) al presentar.",
    )


def fix_intro_slide(slide) -> None:
    if slide.shapes.title:
        slide.shapes.title.text = "Introducción al tema (§1.1)"
    for shape in slide.shapes:
        if shape.has_text_frame and shape != slide.shapes.title:
            tf = shape.text_frame
            tf.clear()
            bullets = [
                "OpenAPI: contrato YAML/JSON machine-readable (Casas et al., 2021)",
                "BIAN: Service Domains = referencia de negocio (Farzi, 2021)",
                "Alineación = correspondencia estructural + semántica contrato ↔ SD",
                "Producto: S, E, C → AlignmentScore (procedimiento reproducible)",
                "Guía 1: diseño documentado · validación empírica → Tesis 2",
            ]
            for i, line in enumerate(bullets):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = line
                p.font.name = FONT
                p.font.size = Pt(22)
            break


def add_caption_box(slide, text: str, top=Inches(6.75)) -> None:
    """Pie de slide para imágenes sin texto (≥20 pt)."""
    existing = [
        s for s in slide.shapes
        if s.has_text_frame and s.text.strip() and s != slide.shapes.title
    ]
    if existing:
        return
    box = slide.shapes.add_textbox(Inches(0.5), top, Inches(12.3), Inches(0.55))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.name = FONT
    p.font.size = Pt(20)
    p.alignment = PP_ALIGN.CENTER


def optimize() -> None:
    if not INPUT.exists():
        raise SystemExit(f"No existe: {INPUT}")

    if not BACKUP.exists():
        import shutil
        shutil.copy2(INPUT, BACKUP)
        print(f"Backup: {BACKUP.name}")

    prs = Presentation(str(INPUT))

    # Eliminar duplicados (índice 0-based, de mayor a menor)
    for n in sorted(DELETE_SLIDES_1BASED, reverse=True):
        delete_slide(prs, n - 1)
        print(f"Eliminada slide duplicada #{n}")

    # Correcciones de contenido
    fix_title_slide(prs.slides[0])
    fix_fuentes_slide(prs.slides[1])
    fix_intro_slide(prs.slides[2])

    # Captions en slides solo-imagen (títulos post-delete)
    captions = {
        3: "Ontología interna — objeto · UA · producto · referencia BIAN",
        4: "Situación real — sector bancario, actores, open banking (CMA, 2016)",
        5: "Síntomas: TSB (2018) · fragmentación · Casas et al. (2021) 68 %",
        6: "Situación subyacente — contrato OpenAPI ↔ Service Domain BIAN",
        7: "Objeto real: taxonomía OpenAPI 3.x · ciclo de vida del contrato",
        8: "Instrumento: parser OpenAPI + extracto BIAN → representaciones normalizadas",
        9: "Casos A (S) y B (C) · clasificación Alta / Media / Baja / Nula",
    }
    for idx, cap in captions.items():
        if idx < len(prs.slides):
            add_caption_box(prs.slides[idx], cap)

    # Notas orales 15 min
    for i, slide in enumerate(prs.slides):
        if i < len(TIMING):
            start, end, who, note = TIMING[i]
            header = f"[{start}–{end}] {who}" if start != "—" else f"[NO EXPONER] {who}"
            set_notes(slide, header, note)

    prs.save(str(OUTPUT))
    print(f"\nOptimizado: {OUTPUT.name} ({len(prs.slides)} diapositivas)")
    print("Exposición oral: 23 slides (omitir slide 2 Fuentes) ≈ 15 min")


if __name__ == "__main__":
    optimize()
