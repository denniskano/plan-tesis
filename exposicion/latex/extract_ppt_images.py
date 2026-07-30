#!/usr/bin/env python3
"""Extrae imágenes del PPT oficial para la presentación Beamer."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

ROOT = Path(__file__).resolve().parent
PPT = ROOT.parent / "ppt" / "Exposicion_Plan_Tesis_Grupo9.2.pptx"
OUT = ROOT / "ppt-images"

# Slides con figuras en Exposicion_Plan_Tesis_Grupo9.2.pptx
SLIDES = [*range(4, 12), 13, 17, 20]


def main() -> None:
    if not PPT.exists():
        raise SystemExit(f"No se encontró: {PPT}")

    OUT.mkdir(parents=True, exist_ok=True)
    for old in OUT.glob("*"):
        if old.is_file():
            old.unlink()

    prs = Presentation(str(PPT))
    count = 0

    for sn in SLIDES:
        slide = prs.slides[sn - 1]
        pic_num = 0
        for shape in slide.shapes:
            if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
            pic_num += 1
            ext = shape.image.ext or "png"
            fname = OUT / f"slide{sn:02d}_img{pic_num:02d}.{ext}"
            fname.write_bytes(shape.image.blob)
            count += 1

    print(f"Extraídas {count} imágenes → {OUT}")


if __name__ == "__main__":
    main()
